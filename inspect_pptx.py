from pptx import Presentation
prs = Presentation('output/enterprise_ai.pptx')
print(f'Slides: {len(prs.slides)}')
print(f'Size: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
print()
for i, slide in enumerate(prs.slides, 1):
    shapes = len(slide.shapes)
    texts = [
        s.text_frame.text[:45].strip().replace('\n', ' ')
        for s in slide.shapes
        if hasattr(s, 'text_frame') and s.text_frame.text.strip()
    ]
    title = texts[0] if texts else '(no text)'
    print(f'  Slide {i:2d}: {shapes:2d} shapes | {title}')
print()
print('OK - PPTX looks structurally valid.')
