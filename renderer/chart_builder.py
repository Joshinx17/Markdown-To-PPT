"""
chart_builder.py
────────────────
Generates programmatic, non-copyrighted charts using python-pptx's native
chart engine (backed by embedded Excel data).  All charts inherit the
Meridian colour palette from theme.py.

Supported types:
  • Bar (clustered column)
  • Pie
  • Line
  • Area
"""

from __future__ import annotations

import logging
from typing import List, Optional

from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_LABEL_POSITION,
)
from pptx.oxml.ns import qn
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches, Pt
from pptx.slide import Slide

from renderer.theme import (
    CHART_COLORS, FONT_BODY, FONT_HEADING,
    TEXT_DARK, TEXT_MUTED, WHITE, NAVY,
)

logger = logging.getLogger(__name__)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _set_series_color(chart, series_idx: int, color: RGBColor) -> None:
    """Apply a solid fill colour to a chart series via XML."""
    try:
        ser = chart.series[series_idx]
        xml_ser = ser._element

        # Build <c:spPr><a:solidFill><a:srgbClr>
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        spPr = xml_ser.find(qn('c:spPr'))
        if spPr is None:
            from lxml import etree
            spPr = etree.SubElement(xml_ser, qn('c:spPr'))

        # Remove existing fills
        for old in spPr.findall(qn('a:solidFill')):
            spPr.remove(old)
        for old in spPr.findall(qn('a:gradFill')):
            spPr.remove(old)

        from lxml import etree
        solid = etree.SubElement(spPr, qn('a:solidFill'))
        srgb  = etree.SubElement(solid, qn('a:srgbClr'))
        srgb.set('val', f'{color.rgb:06X}')
    except Exception as exc:
        logger.debug('Could not colour series %d: %s', series_idx, exc)


def _style_chart(chart, title: str, show_legend: bool = True) -> None:
    """Apply common styling: title, font, no gridline clutter."""
    # Chart title
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        tf = chart.chart_title.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size   = Pt(13)
                run.font.bold   = True
                run.font.name   = FONT_HEADING
                run.font.color.rgb = TEXT_DARK

    # Legend
    chart.has_legend = show_legend
    if show_legend and chart.has_legend:
        chart.legend.position       = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.size  = Pt(11)
            chart.legend.font.name  = FONT_BODY
        except Exception:
            pass

    # Value axis font (where applicable)
    try:
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.name = FONT_BODY
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    except Exception:
        pass

    try:
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.name = FONT_BODY
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass


# ─── Bar chart ───────────────────────────────────────────────────────────────

def add_bar_chart(
    slide: Slide,
    x, y, w, h,
    categories: List[str],
    series_data: List[tuple],   # [(name, [values]), ...]
    title: str = '',
    caption: str = '',
) -> Optional[GraphicFrame]:
    """Add a clustered column (bar) chart."""
    if not categories or not series_data:
        logger.warning('add_bar_chart: no data provided, skipping')
        return None
    try:
        cd = ChartData()
        cd.categories = categories
        for name, values in series_data:
            cd.add_series(name, [float(v) for v in values])

        gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
        chart = gf.chart

        _style_chart(chart, title, show_legend=len(series_data) > 1)

        # colour each series
        for i, _ in enumerate(series_data):
            _set_series_color(chart, i, CHART_COLORS[i % len(CHART_COLORS)])

        # Plot area background → white
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = WHITE
        except Exception:
            pass

        # Caption text box
        if caption:
            _add_caption(slide, caption, x, y + h, w)

        return gf
    except Exception as exc:
        logger.error('add_bar_chart failed: %s', exc)
        return None


# ─── Pie chart ───────────────────────────────────────────────────────────────

def add_pie_chart(
    slide: Slide,
    x, y, w, h,
    categories: List[str],
    values: List[float],
    title: str = '',
    caption: str = '',
) -> Optional[GraphicFrame]:
    """Add a pie chart with data labels."""
    if not categories or not values:
        logger.warning('add_pie_chart: no data provided, skipping')
        return None
    try:
        cd = ChartData()
        cd.categories = categories
        cd.add_series('', [float(v) for v in values])

        gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, cd)
        chart = gf.chart

        _style_chart(chart, title, show_legend=True)

        # Data labels as percentages
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage  = True
            dl.show_value       = False
            dl.show_category_name = True
            dl.font.size        = Pt(10)
            dl.font.name        = FONT_BODY
        except Exception:
            pass

        # Colour pie slices
        try:
            ser = chart.series[0]
            for j, pt in enumerate(ser.points):
                fill = pt.format.fill
                fill.solid()
                fill.fore_color.rgb = CHART_COLORS[j % len(CHART_COLORS)]
        except Exception:
            pass

        if caption:
            _add_caption(slide, caption, x, y + h, w)

        return gf
    except Exception as exc:
        logger.error('add_pie_chart failed: %s', exc)
        return None


# ─── Line chart ──────────────────────────────────────────────────────────────

def add_line_chart(
    slide: Slide,
    x, y, w, h,
    categories: List[str],
    series_data: List[tuple],
    title: str = '',
    caption: str = '',
) -> Optional[GraphicFrame]:
    """Add a line chart."""
    if not categories or not series_data:
        return None
    try:
        cd = ChartData()
        cd.categories = categories
        for name, values in series_data:
            cd.add_series(name, [float(v) for v in values])

        gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, cd)
        chart = gf.chart

        _style_chart(chart, title, show_legend=len(series_data) > 1)

        for i, _ in enumerate(series_data):
            _set_series_color(chart, i, CHART_COLORS[i % len(CHART_COLORS)])

        if caption:
            _add_caption(slide, caption, x, y + h, w)

        return gf
    except Exception as exc:
        logger.error('add_line_chart failed: %s', exc)
        return None


# ─── Area chart ──────────────────────────────────────────────────────────────

def add_area_chart(
    slide: Slide,
    x, y, w, h,
    categories: List[str],
    series_data: List[tuple],
    title: str = '',
    caption: str = '',
) -> Optional[GraphicFrame]:
    """Add a stacked area chart."""
    if not categories or not series_data:
        return None
    try:
        cd = ChartData()
        cd.categories = categories
        for name, values in series_data:
            cd.add_series(name, [float(v) for v in values])

        gf = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x, y, w, h, cd)
        chart = gf.chart

        _style_chart(chart, title, show_legend=len(series_data) > 1)

        for i, _ in enumerate(series_data):
            _set_series_color(chart, i, CHART_COLORS[i % len(CHART_COLORS)])

        if caption:
            _add_caption(slide, caption, x, y + h, w)

        return gf
    except Exception as exc:
        logger.error('add_area_chart failed: %s', exc)
        return None


# ─── Utility: caption text ────────────────────────────────────────────────────

def _add_caption(slide: Slide, caption: str, x, y, w) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    txb = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size  = Pt(10)
    run.font.name  = FONT_BODY
    run.font.italic = True
    run.font.color.rgb = TEXT_MUTED
