"""
infographic_builder.py
──────────────────────
Programmatically generates infographic elements using python-pptx shapes.
All elements inherit the Meridian colour palette and require NO external
images or copyrighted assets.

Infographic types:
  • Process flow  – horizontal arrow-connected step boxes (3-6 steps)
  • Timeline      – horizontal timeline bar with alternating event cards
  • Comparison    – side-by-side labelled columns with bullet points

Design principles:
  • Rounded rectangles for boxes
  • Filled circles for markers / numbers
  • Thin connector lines / arrow shapes
  • Consistent typography from theme.py
"""

from __future__ import annotations

import logging
from typing import List

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.slide import Slide

from renderer.theme import (
    NAVY, NAVY_LIGHT, GOLD, GOLD_LIGHT,
    WHITE, CARD_BG, TEXT_DARK, TEXT_BODY, TEXT_MUTED,
    FONT_HEADING, FONT_BODY, CHART_COLORS,
)

logger = logging.getLogger(__name__)

# ─── Internal shape helpers ───────────────────────────────────────────────────

# MSO AutoShape IDs (match VBA MsoAutoShapeType enum)
_ROUNDED_RECT  = 5     # msoShapeRoundedRectangle
_OVAL          = 9     # msoShapeOval
_RECT          = 1     # msoShapeRectangle
_RIGHT_ARROW   = 13    # msoShapeRightArrow
_CHEVRON       = 55    # msoShapeChevron  (may not be available everywhere)


def _add_shape(slide: Slide, shape_id: int, x, y, w, h,
               fill: RGBColor, line_color: RGBColor = None,
               line_width: float = 0.0) -> object:
    """Add an auto-shape with solid fill and optional border."""
    shape = slide.shapes.add_shape(shape_id, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()  # no border
    return shape


def _add_text_box(slide: Slide, text: str, x, y, w, h,
                  font_size: int, bold: bool, color: RGBColor,
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  font_name: str = FONT_BODY,
                  wrap: bool = True) -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(font_size)
    run.font.bold       = bold
    run.font.name       = font_name
    run.font.color.rgb  = color


def _add_circle_number(slide: Slide, number: str, cx, cy, diameter,
                       bg: RGBColor, text_color: RGBColor) -> None:
    """Draw a filled circle centred at (cx,cy) with a number label inside."""
    x = cx - diameter // 2
    y = cy - diameter // 2
    shape = _add_shape(slide, _OVAL, x, y, diameter, diameter, fill=bg)

    tf = shape.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.auto_size = None
    # vertical centering
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    run = p.add_run()
    run.text           = str(number)
    run.font.size      = Pt(max(8, int(diameter / Inches(0.05)) // 2))  # heuristic
    run.font.bold      = True
    run.font.name      = FONT_HEADING
    run.font.color.rgb = text_color


# ─── Process Flow ─────────────────────────────────────────────────────────────

def build_process_flow(
    slide: Slide,
    steps: List[object],    # List[ProcessStep]
    x, y, w, h,
) -> None:
    """
    Draw a left-to-right process flow diagram with numbered step boxes
    and gold arrow connectors between them.

    Handles 2-6 steps.  For > 6 steps uses a two-row layout.
    """
    n = len(steps)
    if n == 0:
        return
    n = min(n, 6)
    steps = steps[:n]

    # ── Layout decision ──────────────────────────────────────────────────────
    if n <= 5:
        _process_flow_single_row(slide, steps, x, y, w, h)
    else:
        _process_flow_double_row(slide, steps, x, y, w, h)


def _process_flow_single_row(slide, steps, x, y, w, h):
    n = len(steps)
    arrow_w  = Inches(0.35)
    gap      = arrow_w
    box_w    = (w - gap * (n - 1)) / n
    box_h    = Inches(1.60)
    num_d    = Inches(0.42)

    box_y    = y + (h - box_h - num_d * 0.6) / 2 + num_d * 0.6
    num_y    = box_y - num_d * 0.75

    for i, step in enumerate(steps):
        bx = x + i * (box_w + gap)
        color = CHART_COLORS[i % len(CHART_COLORS)]

        # Number circle
        cx = int(bx + box_w / 2)
        cy = int(num_y + num_d / 2)
        _add_circle_number(slide, str(i + 1), cx, cy, num_d, color, WHITE)

        # Step box
        box = _add_shape(slide, _ROUNDED_RECT, bx, box_y, box_w, box_h,
                         fill=color)
        tf = box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text           = step.title[:35]
        run.font.size      = Pt(13)
        run.font.bold      = True
        run.font.name      = FONT_HEADING
        run.font.color.rgb = WHITE

        # Description
        if step.description:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text           = step.description[:60]
            r2.font.size      = Pt(10)
            r2.font.name      = FONT_BODY
            r2.font.color.rgb = RGBColor(0xDD, 0xE8, 0xFF)

        # Arrow → (except after last box)
        if i < n - 1:
            ax = bx + box_w + Inches(0.04)
            ay = box_y + box_h / 2 - Inches(0.13)
            arw = _add_shape(slide, _RIGHT_ARROW,
                             ax, ay, gap - Inches(0.08), Inches(0.26),
                             fill=GOLD)


def _process_flow_double_row(slide, steps, x, y, w, h):
    """Two-row layout for 6 steps."""
    top_steps    = steps[:3]
    bottom_steps = steps[3:]
    row_h = h / 2 - Inches(0.1)
    _process_flow_single_row(slide, top_steps, x, y, w, row_h)
    _process_flow_single_row(slide, bottom_steps, x, y + row_h + Inches(0.15), w, row_h)


# ─── Timeline ─────────────────────────────────────────────────────────────────

def build_timeline(
    slide: Slide,
    events: List[object],   # List[TimelineEvent]
    x, y, w, h,
) -> None:
    """
    Horizontal timeline:  ──●───●───●─── with alternating above/below cards.
    Handles up to 8 events.
    """
    if not events:
        return
    events = events[:8]
    n = len(events)

    line_y   = y + h * 0.48
    line_h   = Inches(0.06)
    # Main line
    line_bar = _add_shape(slide, _RECT,
                          x, int(line_y - line_h / 2), w, line_h,
                          fill=NAVY)

    dot_d    = Inches(0.22)
    card_w   = max(Inches(1.4), (w - Inches(0.5)) / n - Inches(0.1))
    card_h   = Inches(1.20)

    spacing  = w / n

    for i, ev in enumerate(events):
        cx = int(x + i * spacing + spacing / 2)
        cy = int(line_y)

        # Dot on line
        color = CHART_COLORS[i % len(CHART_COLORS)]
        _add_circle_number(slide, str(i + 1), cx, cy, dot_d, color, WHITE)

        # Alternate above / below
        above = (i % 2 == 0)
        card_x = cx - card_w // 2

        label_size = 9
        title_size = 11

        if above:
            # Date label just above card
            date_y = int(line_y - card_h - Inches(0.55) - dot_d / 2)
            _add_text_box(slide, ev.date, card_x, date_y,
                          card_w, Inches(0.25), label_size, True, color,
                          PP_ALIGN.CENTER)
            # Card
            card_y = int(line_y - card_h - Inches(0.15) - dot_d / 2)
        else:
            card_y = int(line_y + dot_d / 2 + Inches(0.15))
            date_y = int(card_y + card_h + Inches(0.04))
            _add_text_box(slide, ev.date, card_x, date_y,
                          card_w, Inches(0.25), label_size, True, color,
                          PP_ALIGN.CENTER)

        # Card box
        card = _add_shape(slide, _ROUNDED_RECT,
                          card_x, card_y, card_w, card_h,
                          fill=CARD_BG, line_color=color, line_width=1.0)
        tf = card.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text           = ev.title[:30]
        r.font.size      = Pt(title_size)
        r.font.bold      = True
        r.font.name      = FONT_HEADING
        r.font.color.rgb = NAVY

        if ev.description:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text           = ev.description[:50]
            r2.font.size      = Pt(9)
            r2.font.name      = FONT_BODY
            r2.font.color.rgb = TEXT_BODY


# ─── Comparison grid ──────────────────────────────────────────────────────────

def build_comparison(
    slide: Slide,
    items: List[object],    # List[ComparisonItem]
    x, y, w, h,
) -> None:
    """
    Side-by-side comparison columns, each with a coloured header and bullets.
    Handles 2-4 items.
    """
    if not items:
        return
    items = items[:4]
    n = len(items)

    col_w    = (w - Inches(0.15) * (n - 1)) / n
    gap      = Inches(0.15)
    hdr_h    = Inches(0.55)

    for i, item in enumerate(items):
        cx = x + i * (col_w + gap)
        color = CHART_COLORS[i % len(CHART_COLORS)]

        # Header box
        hdr = _add_shape(slide, _ROUNDED_RECT,
                         cx, y, col_w, hdr_h, fill=color)
        tf = hdr.text_frame
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text           = item.label[:30]
        r.font.size      = Pt(14)
        r.font.bold      = True
        r.font.name      = FONT_HEADING
        r.font.color.rgb = WHITE

        # Body card
        body_y = y + hdr_h + Inches(0.08)
        body_h = h - hdr_h - Inches(0.08)
        body = _add_shape(slide, _ROUNDED_RECT,
                          cx, body_y, col_w, body_h,
                          fill=CARD_BG, line_color=color, line_width=0.75)

        bf = body.text_frame
        bf.word_wrap = True
        bf.margin_left   = Inches(0.12)
        bf.margin_right  = Inches(0.08)
        bf.margin_top    = Inches(0.10)
        bf.margin_bottom = Inches(0.10)

        first = True
        for point in item.points[:6]:
            p2 = bf.paragraphs[0] if first else bf.add_paragraph()
            first = False
            r2 = p2.add_run()
            r2.text           = f'▶  {point[:55]}'
            r2.font.size      = Pt(12)
            r2.font.name      = FONT_BODY
            r2.font.color.rgb = TEXT_BODY
