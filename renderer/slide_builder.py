"""
slide_builder.py
────────────────
Core rendering engine: converts each SlideBlueprint into a populated
python-pptx Slide object.

Slide types handled:
  TITLE · AGENDA · EXECUTIVE_SUMMARY · SECTION_DIVIDER · CONTENT_BULLETS
  CONTENT_TEXT · TWO_COLUMN · CHART_BAR · CHART_PIE · CHART_LINE
  CHART_AREA · TABLE · PROCESS_FLOW · TIMELINE · COMPARISON · CONCLUSION
  QUOTE

Modern Design Language: "Meridian Professional"
  • Navy headers with teal accent stripe for modern appeal
  • Professional spacing and typography with Segoe UI
  • Rounded corners, professional shadows, modern color palette
  • Maximum 5 bullets per slide, 15 words per bullet
  • No external images; all decorative elements are geometric shapes
"""

from __future__ import annotations

import logging
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.slide import Slide

from renderer import theme as T
from renderer.chart_builder import (
    add_bar_chart, add_pie_chart, add_line_chart, add_area_chart,
)
from renderer.table_builder import add_table
from renderer.infographic_builder import (
    build_process_flow, build_timeline, build_comparison,
    _add_shape, _add_text_box, _ROUNDED_RECT, _OVAL, _RECT,
)
from structurer.slide_types import (
    PresentationBlueprint, SlideBlueprint, SlideType,
)

logger = logging.getLogger(__name__)


# ─── Presentation factory ─────────────────────────────────────────────────────

def create_presentation() -> Presentation:
    """Return a blank widescreen 16:9 Presentation."""
    prs = Presentation()
    prs.slide_width  = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _blank_slide(prs: Presentation) -> Slide:
    """Add a blank (layout 6) slide and return it."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank in default master
    return prs.slides.add_slide(blank_layout)


# ─── Common structural elements ───────────────────────────────────────────────

def _set_background(slide: Slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header(slide: Slide, title: str, slide_num: int = 0, total: int = 0) -> None:
    """Modern navy header bar with teal accent stripe and slide title."""
    # Navy background bar
    bar = _add_shape(slide, _RECT,
                     0, 0, T.SLIDE_W, T.HEADER_H, fill=T.NAVY_DARK)

    # Teal left accent stripe (modern professional accent)
    _add_shape(slide, _RECT,
               0, 0, T.ACCENT_BAR_W, T.HEADER_H, fill=T.TEAL)

    # Title text (right of accent stripe)
    txb = slide.shapes.add_textbox(
        T.ACCENT_BAR_W + Inches(0.25), Inches(0.18),
        T.SLIDE_W - T.ACCENT_BAR_W - Inches(0.60), T.HEADER_H - Inches(0.40),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title[:85]
    run.font.size      = T.SIZE_SLIDE_TITLE
    run.font.bold      = True
    run.font.name      = T.FONT_HEADING
    run.font.color.rgb = T.WHITE


def _add_footer(slide: Slide, slide_num: int, total: int,
                deck_title: str = '') -> None:
    """Thin footer bar with deck title (left) and page number (right)."""
    # Background
    _add_shape(slide, _RECT,
               0, T.FOOTER_Y, T.SLIDE_W, T.FOOTER_H,
               fill=T.NAVY)

    # Deck title (left)
    if deck_title:
        txb = slide.shapes.add_textbox(
            Inches(0.25), T.FOOTER_Y, T.SLIDE_W * 0.7, T.FOOTER_H,
        )
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text           = deck_title[:60]
        r.font.size      = T.SIZE_FOOTER
        r.font.name      = T.FONT_BODY
        r.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)

    # Slide number (right)
    num_text = f'{slide_num} / {total}' if total else str(slide_num)
    txb2 = slide.shapes.add_textbox(
        T.SLIDE_W - Inches(1.2), T.FOOTER_Y, Inches(1.1), T.FOOTER_H,
    )
    tf2 = txb2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2  = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text           = num_text
    r2.font.size      = T.SIZE_FOOTER
    r2.font.name      = T.FONT_BODY
    r2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)


def _add_speaker_notes(slide: Slide, notes_text: str) -> None:
    if not notes_text:
        return
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = notes_text
    except Exception:
        pass


# ─── Content helpers ──────────────────────────────────────────────────────────

def _add_bullet_list(slide: Slide, bullets, x, y, w, h) -> None:
    """Render a list of BulletItem objects as styled text with teal accents."""
    if not bullets:
        return

    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True

    # Modern bullet indicator colours by level (teal / navy gradient)
    indicator_colors = [T.TEAL, T.NAVY_LIGHT, T.TEXT_MUTED]
    indicator_chars = ['▸', '–', '◦']  # modern bullet chars

    first = True
    for i, b in enumerate(bullets[:5]):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        # Indicator (teal bullet or dash)
        indicator = indicator_chars[min(b.level, 2)]
        ind_color  = indicator_colors[min(b.level, 2)]

        r_ind = p.add_run()
        r_ind.text           = f'{indicator}  '
        r_ind.font.size      = Pt(16 - b.level * 2)
        r_ind.font.bold      = b.level == 0
        r_ind.font.name      = T.FONT_BODY
        r_ind.font.color.rgb = ind_color

        # Main text
        r_text = p.add_run()
        r_text.text           = b.text[:120]
        r_text.font.size      = T.SIZE_BODY if b.level == 0 else Pt(15)
        r_text.font.bold      = b.level == 0
        r_text.font.name      = T.FONT_BODY
        r_text.font.color.rgb = T.TEXT_DARK if b.level == 0 else T.TEXT_BODY

        # Modern spacing
        p.space_before = Pt(10 if b.level == 0 else 5)
        p.space_after  = Pt(3)
        p.line_spacing = 1.15  # improve readability


def _add_callout_card(slide: Slide, icon: str, title: str, desc: str,
                       x, y, w, h, accent_color: RGBColor) -> None:
    """Rounded card with coloured top-bar accent, icon circle, title, body."""
    # Card background
    card = _add_shape(slide, _ROUNDED_RECT, x, y, w, h, fill=T.CARD_BG,
                      line_color=accent_color, line_width=1.2)

    # Top accent bar
    _add_shape(slide, _RECT, x, y, w, Inches(0.32), fill=accent_color)

    # Icon in circle
    icon_d  = Inches(0.52)
    icon_cx = int(x + w // 2)
    icon_cy = int(y + Inches(0.32) + icon_d // 2 + Inches(0.08))
    circ = _add_shape(slide, _OVAL,
                      icon_cx - icon_d // 2, icon_cy - icon_d // 2,
                      icon_d, icon_d,
                      fill=accent_color)
    icon_tf = circ.text_frame
    icon_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ip = icon_tf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    ir = ip.add_run()
    ir.text = icon[:2]
    ir.font.size = Pt(18)

    # Title
    title_y = icon_cy + icon_d // 2 + Inches(0.10)
    title_h = Inches(0.45)
    _add_text_box(slide, title[:40], x + Inches(0.12), title_y,
                  w - Inches(0.24), title_h,
                  font_size=14, bold=True, color=T.NAVY,
                  align=PP_ALIGN.CENTER)

    # Description
    desc_y = title_y + title_h
    desc_h = y + h - desc_y - Inches(0.10)
    _add_text_box(slide, desc[:160], x + Inches(0.12), desc_y,
                  w - Inches(0.24), desc_h,
                  font_size=11, bold=False, color=T.TEXT_BODY,
                  align=PP_ALIGN.CENTER)


def _add_bullet_cards(slide: Slide, bullets, x, y, w, h) -> bool:
    top_level = [bullet for bullet in bullets if bullet.level == 0][:4]
    if len(top_level) < 2:
        return False

    cols = 2 if len(top_level) > 2 else 1
    rows = (len(top_level) + cols - 1) // cols
    gap_x = Inches(0.22)
    gap_y = Inches(0.18)
    card_w = (w - gap_x * (cols - 1)) / cols
    card_h = (h - gap_y * (rows - 1)) / rows

    for index, bullet in enumerate(top_level):
        row = index // cols
        col = index % cols
        cx = x + col * (card_w + gap_x)
        cy = y + row * (card_h + gap_y)
        accent = T.CHART_COLORS[index % len(T.CHART_COLORS)]

        _add_shape(slide, _ROUNDED_RECT, cx, cy, card_w, card_h, fill=T.CARD_BG, line_color=accent, line_width=1.1)
        _add_shape(slide, _RECT, cx, cy, card_w, Inches(0.10), fill=accent)

        badge = _add_shape(slide, _OVAL, cx + Inches(0.18), cy + Inches(0.20), Inches(0.46), Inches(0.46), fill=accent)
        badge_tf = badge.text_frame
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        badge_p = badge_tf.paragraphs[0]
        badge_p.alignment = PP_ALIGN.CENTER
        badge_r = badge_p.add_run()
        badge_r.text = str(index + 1)
        badge_r.font.size = Pt(14)
        badge_r.font.bold = True
        badge_r.font.name = T.FONT_HEADING
        badge_r.font.color.rgb = T.WHITE

        text = bullet.text[:180]
        _add_text_box(
            slide,
            text,
            cx + Inches(0.82),
            cy + Inches(0.16),
            card_w - Inches(1.00),
            card_h - Inches(0.28),
            font_size=16,
            bold=True,
            color=T.TEXT_DARK,
            align=PP_ALIGN.LEFT,
        )
    return True


# ─── Slide builders (one per SlideType) ──────────────────────────────────────

def _build_title(slide: Slide, bp: SlideBlueprint, deck_title: str) -> None:
    _set_background(slide, T.NAVY)

    # Decorative shapes: semi-transparent lighter navy rectangles
    _add_shape(slide, _RECT,
               T.SLIDE_W - Inches(3.5), T.SLIDE_H - Inches(3.5),
               Inches(3.5), Inches(3.5),
               fill=T.NAVY_MID)
    _add_shape(slide, _RECT,
               T.SLIDE_W - Inches(2.0), T.SLIDE_H - Inches(2.0),
               Inches(2.0), Inches(2.0),
               fill=T.NAVY_LIGHT)

    # Teal horizontal accent line (modern professional accent)
    accent_y = Inches(4.50)
    _add_shape(slide, _RECT,
               Inches(0.80), accent_y, Inches(4.5), Inches(0.055),
               fill=T.TEAL)

    # Main title
    txb_title = slide.shapes.add_textbox(
        Inches(0.80), Inches(1.80), Inches(11.0), Inches(2.40),
    )
    ttf = txb_title.text_frame
    ttf.word_wrap = True
    ttf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text           = bp.title[:120]
    tr.font.size      = T.SIZE_TITLE_MAIN
    tr.font.bold      = True
    tr.font.name      = T.FONT_HEADING
    tr.font.color.rgb = T.WHITE

    # Subtitle
    if bp.subtitle:
        txb_sub = slide.shapes.add_textbox(
            Inches(0.80), Inches(4.65), Inches(11.0), Inches(1.00),
        )
        stf = txb_sub.text_frame
        stf.word_wrap = True
        sp  = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr  = sp.add_run()
        sr.text           = bp.subtitle[:200]
        sr.font.size      = T.SIZE_TITLE_SUB
        sr.font.name      = T.FONT_BODY
        sr.font.color.rgb = T.TEAL_LIGHT

    # Author / date bottom
    meta_parts = []
    if bp.author:
        meta_parts.append(bp.author)
    if bp.date:
        meta_parts.append(bp.date)
    if meta_parts:
        txb_meta = slide.shapes.add_textbox(
            Inches(0.80), Inches(6.60), Inches(8.0), Inches(0.60),
        )
        mf = txb_meta.text_frame
        mp = mf.paragraphs[0]
        mp.alignment = PP_ALIGN.LEFT
        mr  = mp.add_run()
        mr.text = '  ·  '.join(meta_parts)
        mr.font.size      = Pt(13)
        mr.font.name      = T.FONT_BODY
        mr.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)


def _build_agenda(slide: Slide, bp: SlideBlueprint, deck_title: str,
                  slide_num: int, total: int) -> None:
    _set_background(slide, T.OFF_WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    # Two-column layout for agenda items
    items = bp.bullets[:8]
    col_count = 2 if len(items) > 4 else 1
    col_w = T.CONTENT_W / col_count - Inches(0.15)
    col_h = T.CONTENT_H
    col_items_each = (len(items) + col_count - 1) // col_count

    for col_idx in range(col_count):
        col_x = T.CONTENT_X + col_idx * (col_w + Inches(0.3))
        col_items = items[col_idx * col_items_each: (col_idx + 1) * col_items_each]

        ty = T.CONTENT_Y + Inches(0.15)
        item_h = Inches(0.68)
        num_d  = Inches(0.46)

        for i, bullet in enumerate(col_items):
            cy = int(ty + num_d / 2)
            # Number circle
            circ = _add_shape(slide, _OVAL,
                              col_x, int(ty), num_d, num_d,
                              fill=T.NAVY)
            ctf = circ.text_frame
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            # Global agenda numbering
            global_idx = col_idx * col_items_each + i + 1
            cr.text = str(global_idx)
            cr.font.size      = Pt(13)
            cr.font.bold      = True
            cr.font.name      = T.FONT_HEADING
            cr.font.color.rgb = T.WHITE

            # Item text
            txt_x = col_x + num_d + Inches(0.15)
            txt_w = col_w - num_d - Inches(0.15)
            _add_text_box(slide, bullet.text[:80],
                          txt_x, int(ty + Inches(0.05)),
                          txt_w, num_d - Inches(0.05),
                          font_size=17, bold=False, color=T.TEXT_DARK,
                          align=PP_ALIGN.LEFT)
            ty += item_h


def _build_executive_summary(slide: Slide, bp: SlideBlueprint, deck_title: str,
                              slide_num: int, total: int) -> None:
    _set_background(slide, T.OFF_WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    points = bp.summary_points[:4]
    n = len(points)
    if n == 0:
        return

    # Dynamic layout: 2 or 3 or 4 per row
    cards_per_row = min(n, 3)
    rows = (n + cards_per_row - 1) // cards_per_row

    area_w = T.CONTENT_W
    area_h = T.CONTENT_H - Inches(0.15)
    gap    = Inches(0.22)

    card_w = (area_w - gap * (cards_per_row - 1)) / cards_per_row
    card_h = (area_h - gap * (rows - 1)) / rows

    for i, pt in enumerate(points):
        row = i // cards_per_row
        col = i % cards_per_row
        cx  = T.CONTENT_X + col * (card_w + gap)
        cy  = T.CONTENT_Y + Inches(0.10) + row * (card_h + gap)
        accent = T.CHART_COLORS[i % len(T.CHART_COLORS)]
        _add_callout_card(slide, pt.icon, pt.title, pt.description,
                          cx, cy, card_w, card_h, accent)


def _build_section_divider(slide: Slide, bp: SlideBlueprint, deck_title: str,
                            slide_num: int, total: int) -> None:
    _set_background(slide, T.NAVY_MID)
    _add_footer(slide, slide_num, total, deck_title)

    # Decorative teal elements (modern professional design)
    _add_shape(slide, _RECT, 0, 0, Inches(0.40), T.SLIDE_H, fill=T.TEAL)
    _add_shape(slide, _OVAL, T.SLIDE_W - Inches(2.8), Inches(0.8), Inches(2.0), Inches(2.0), fill=T.NAVY_LIGHT)
    _add_shape(slide, _OVAL, T.SLIDE_W - Inches(1.8), Inches(1.8), Inches(1.1), Inches(1.1), fill=T.TEAL)
    _add_shape(slide, _RECT, Inches(0.95), Inches(3.95), Inches(3.6), Inches(0.08), fill=T.TEAL)

    _add_text_box(slide, f'SECTION {slide_num - 2}',
                  Inches(0.90), Inches(2.20), Inches(9.0), Inches(0.45),
                  font_size=12, bold=True,
                  color=RGBColor(0xAA, 0xBB, 0xDD),
                  align=PP_ALIGN.LEFT)

    # Section title
    txb = slide.shapes.add_textbox(
        Inches(0.90), Inches(2.70), Inches(11.0), Inches(2.20),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text           = bp.title[:80]
    r.font.size      = Pt(38)
    r.font.bold      = True
    r.font.name      = T.FONT_HEADING
    r.font.color.rgb = T.WHITE

    # Subtitle / description
    if bp.subtitle:
        _add_text_box(slide, bp.subtitle[:160],
                      Inches(0.90), Inches(4.90), Inches(9.0), Inches(0.80),
                      font_size=18, bold=False,
                      color=T.TEAL_LIGHT,
                      align=PP_ALIGN.LEFT)


def _build_content_bullets(slide: Slide, bp: SlideBlueprint, deck_title: str,
                            slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    # Optional: if body_text exists, show it as a paragraph first
    content_y = T.CONTENT_Y + Inches(0.10)
    if bp.body_text and not bp.bullets:
        txb = slide.shapes.add_textbox(
            T.CONTENT_X + Inches(0.20), content_y,
            T.CONTENT_W - Inches(0.20), T.CONTENT_H,
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text           = bp.body_text[:800]
        r.font.size      = T.SIZE_BODY
        r.font.name      = T.FONT_BODY
        r.font.color.rgb = T.TEXT_BODY
        return

    if _add_bullet_cards(
        slide,
        bp.bullets,
        T.CONTENT_X + Inches(0.12),
        content_y,
        T.CONTENT_W - Inches(0.24),
        T.CONTENT_H - Inches(0.10),
    ):
        return

    _add_bullet_list(slide, bp.bullets,
                     T.CONTENT_X + Inches(0.20), content_y,
                     T.CONTENT_W - Inches(0.20), T.CONTENT_H)


def _build_two_column(slide: Slide, bp: SlideBlueprint, deck_title: str,
                      slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    mid_x    = T.SLIDE_W / 2
    col_pad  = Inches(0.25)
    col_w    = mid_x - T.CONTENT_X - col_pad - Inches(0.10)
    col_y    = T.CONTENT_Y + Inches(0.08)
    col_h    = T.CONTENT_H - Inches(0.08)

    for side in ('left', 'right'):
        if side == 'left':
            cx     = T.CONTENT_X
            title  = bp.left_title
            buls   = bp.left_bullets
            acc    = T.NAVY
        else:
            cx    = mid_x + col_pad
            title  = bp.right_title
            buls   = bp.right_bullets
            acc    = T.TEAL

        # Column title
        _add_shape(slide, _RECT,
                   cx, col_y, col_w, Inches(0.08), fill=acc)
        _add_text_box(slide, title[:50],
                      cx, col_y + Inches(0.10), col_w, Inches(0.50),
                      font_size=16, bold=True, color=T.NAVY,
                      align=PP_ALIGN.LEFT)

        # Bullet items
        items_y = col_y + Inches(0.68)
        from structurer.slide_types import BulletItem
        bullet_objs = [BulletItem(text=b, level=0) for b in buls[:5]]
        _add_bullet_list(slide, bullet_objs,
                         cx, items_y, col_w, col_h - Inches(0.68))

    # Divider line
    div_x = int(mid_x - Inches(0.04))
    _add_shape(slide, _RECT,
               div_x, T.CONTENT_Y + Inches(0.08),
               Inches(0.03), T.CONTENT_H - Inches(0.08),
               fill=T.BORDER)


def _build_chart(slide: Slide, bp: SlideBlueprint, deck_title: str,
                 slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    chart_x = T.CONTENT_X + Inches(0.20)
    chart_y = T.CONTENT_Y + Inches(0.20)
    chart_w = T.CONTENT_W - Inches(0.40)
    chart_h = T.CONTENT_H - Inches(0.60)

    series_data = [(s.name, s.values) for s in bp.series]

    if bp.type == SlideType.CHART_BAR:
        add_bar_chart(slide, chart_x, chart_y, chart_w, chart_h,
                      bp.categories, series_data,
                      title=bp.chart_title, caption=bp.chart_caption)

    elif bp.type == SlideType.CHART_PIE:
        add_pie_chart(slide, chart_x, chart_y, chart_w, chart_h,
                      bp.categories, bp.values,
                      title=bp.chart_title, caption=bp.chart_caption)

    elif bp.type == SlideType.CHART_LINE:
        add_line_chart(slide, chart_x, chart_y, chart_w, chart_h,
                       bp.categories, series_data,
                       title=bp.chart_title, caption=bp.chart_caption)

    elif bp.type == SlideType.CHART_AREA:
        add_area_chart(slide, chart_x, chart_y, chart_w, chart_h,
                       bp.categories, series_data,
                       title=bp.chart_title, caption=bp.chart_caption)


def _build_table(slide: Slide, bp: SlideBlueprint, deck_title: str,
                 slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    n_rows = min(len(bp.rows), 14)
    n_cols = len(bp.headers)

    tbl_x = T.CONTENT_X + Inches(0.15)
    tbl_y = T.CONTENT_Y + Inches(0.20)
    tbl_w = T.CONTENT_W - Inches(0.30)

    # Auto height: ~0.36" per row + header
    row_h_each = Inches(0.38)
    tbl_h = min((n_rows + 1) * row_h_each + Inches(0.10),
                T.CONTENT_H - Inches(0.40))

    add_table(slide, tbl_x, tbl_y, tbl_w, tbl_h,
              bp.headers, bp.rows[:n_rows])


def _build_process_flow(slide: Slide, bp: SlideBlueprint, deck_title: str,
                        slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    build_process_flow(
        slide, bp.steps,
        x=T.CONTENT_X, y=T.CONTENT_Y + Inches(0.20),
        w=T.CONTENT_W, h=T.CONTENT_H - Inches(0.20),
    )


def _build_timeline(slide: Slide, bp: SlideBlueprint, deck_title: str,
                    slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    build_timeline(
        slide, bp.events,
        x=T.CONTENT_X, y=T.CONTENT_Y + Inches(0.10),
        w=T.CONTENT_W, h=T.CONTENT_H - Inches(0.10),
    )


def _build_comparison(slide: Slide, bp: SlideBlueprint, deck_title: str,
                      slide_num: int, total: int) -> None:
    _set_background(slide, T.WHITE)
    _add_header(slide, bp.title, slide_num, total)
    _add_footer(slide, slide_num, total, deck_title)

    build_comparison(
        slide, bp.comparison_items,
        x=T.CONTENT_X, y=T.CONTENT_Y + Inches(0.15),
        w=T.CONTENT_W, h=T.CONTENT_H - Inches(0.15),
    )


def _build_conclusion(slide: Slide, bp: SlideBlueprint, deck_title: str,
                      slide_num: int, total: int) -> None:
    _set_background(slide, T.NAVY)

    # Teal accent bar top (modern professional design)
    _add_shape(slide, _RECT, 0, 0, T.SLIDE_W, Inches(0.12), fill=T.TEAL)

    # Title
    _add_text_box(slide, bp.title,
                  Inches(0.80), Inches(0.20), Inches(10.0), Inches(0.90),
                  font_size=30, bold=True, color=T.WHITE,
                  align=PP_ALIGN.LEFT, font_name=T.FONT_HEADING)

    # Footer
    _add_footer(slide, slide_num, total, deck_title)

    takeaways = bp.takeaways[:6]
    n = len(takeaways)
    if n == 0:
        return

    cols = 2 if n > 3 else 1
    rows = (n + cols - 1) // cols

    area_x = Inches(0.50)
    area_y = Inches(1.20)
    area_w = T.SLIDE_W - Inches(1.0)
    area_h = T.FOOTER_Y - area_y - Inches(0.10)

    gap_x = Inches(0.20)
    gap_y = Inches(0.15)
    card_w = (area_w - gap_x * (cols - 1)) / cols
    card_h = (area_h - gap_y * (rows - 1)) / rows

    for i, tk in enumerate(takeaways):
        row = i // cols
        col = i % cols
        cx  = area_x + col * (card_w + gap_x)
        cy  = area_y + row * (card_h + gap_y)

        # Card bg
        card = _add_shape(slide, _ROUNDED_RECT,
                          cx, cy, card_w, card_h,
                          fill=T.NAVY_MID)

        # Number circle
        num_d = Inches(0.46)
        circ  = _add_shape(slide, _OVAL,
                           cx + Inches(0.15), cy + Inches(0.12),
                           num_d, num_d, fill=T.TEAL)
        ctf = circ.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp  = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr  = cp.add_run()
        cr.text = str(tk.number)
        cr.font.size      = Pt(15)
        cr.font.bold      = True
        cr.font.name      = T.FONT_HEADING
        cr.font.color.rgb = T.NAVY

        # Title + description
        text_x = cx + num_d + Inches(0.30)
        text_w = card_w - num_d - Inches(0.45)
        text_h = card_h - Inches(0.15)
        text_y = cy + Inches(0.10)

        txb = slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
        tf  = txb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text           = tk.title[:60]
        r.font.size      = Pt(14)
        r.font.bold      = True
        r.font.name      = T.FONT_HEADING
        r.font.color.rgb = T.WHITE

        if tk.description:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text           = tk.description[:120]
            r2.font.size      = Pt(11)
            r2.font.name      = T.FONT_BODY
            r2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)


def _build_quote(slide: Slide, bp: SlideBlueprint, deck_title: str,
                 slide_num: int, total: int) -> None:
    _set_background(slide, T.NAVY_MID)
    _add_footer(slide, slide_num, total, deck_title)

    # Large quotation mark decoration
    _add_text_box(slide, '\u201c',
                  Inches(0.60), Inches(0.50), Inches(2.0), Inches(2.0),
                  font_size=96, bold=True, color=T.TEAL,
                  align=PP_ALIGN.LEFT, font_name=T.FONT_HEADING)

    _add_text_box(slide, bp.body_text[:300] or bp.title,
                  Inches(1.20), Inches(1.60), Inches(11.0), Inches(3.50),
                  font_size=24, bold=False, color=T.WHITE,
                  align=PP_ALIGN.LEFT, font_name=T.FONT_BODY)


# ─── Dispatcher ───────────────────────────────────────────────────────────────

_BUILDERS = {
    SlideType.TITLE:             _build_title,
    SlideType.AGENDA:            _build_agenda,
    SlideType.EXECUTIVE_SUMMARY: _build_executive_summary,
    SlideType.SECTION_DIVIDER:   _build_section_divider,
    SlideType.CONTENT_BULLETS:   _build_content_bullets,
    SlideType.CONTENT_TEXT:      _build_content_bullets,   # same renderer
    SlideType.TWO_COLUMN:        _build_two_column,
    SlideType.CHART_BAR:         _build_chart,
    SlideType.CHART_PIE:         _build_chart,
    SlideType.CHART_LINE:        _build_chart,
    SlideType.CHART_AREA:        _build_chart,
    SlideType.TABLE:             _build_table,
    SlideType.PROCESS_FLOW:      _build_process_flow,
    SlideType.TIMELINE:          _build_timeline,
    SlideType.COMPARISON:        _build_comparison,
    SlideType.CONCLUSION:        _build_conclusion,
    SlideType.QUOTE:             _build_quote,
}


def render_presentation(blueprint: PresentationBlueprint) -> Presentation:
    """
    Convert a PresentationBlueprint into a fully rendered python-pptx
    Presentation object.

    Args:
        blueprint: the structured slide plan produced by llm_structurer

    Returns:
        A python-pptx Presentation ready to be saved as .pptx
    """
    prs = create_presentation()
    deck_title = blueprint.presentation_title
    total = blueprint.total_slides

    for bp in blueprint.slides:
        slide = _blank_slide(prs)
        builder = _BUILDERS.get(bp.type)

        if builder is None:
            logger.warning('No builder for %s — using CONTENT_BULLETS', bp.type)
            builder = _build_content_bullets

        try:
            if bp.type == SlideType.TITLE:
                builder(slide, bp, deck_title)
            else:
                builder(slide, bp, deck_title, bp.slide_number, total)
        except Exception as exc:
            logger.error('Error building slide %d (%s): %s',
                         bp.slide_number, bp.type, exc, exc_info=True)
            # Fallback: plain text slide
            try:
                _set_background(slide, T.WHITE)
                _add_header(slide, bp.title, bp.slide_number, total)
                _add_footer(slide, bp.slide_number, total, deck_title)
            except Exception:
                pass

        _add_speaker_notes(slide, bp.speaker_notes)

    logger.info('Rendered %d slides for "%s"', len(blueprint.slides), deck_title)
    return prs
