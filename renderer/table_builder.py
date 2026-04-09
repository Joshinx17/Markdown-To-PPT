"""
table_builder.py
────────────────
Renders styled data tables on slides using python-pptx, with manual cell
colouring and border styling to match the Meridian theme.

Features:
  • Navy header row with white bold text
  • Alternating light-blue / white row colours
  • Centred column headers, left-aligned data cells
  • Auto font-size: shrinks for wide/tall tables
  • Graceful fallback for empty data
"""

from __future__ import annotations

import logging
from typing import List

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.slide import Slide

from renderer.theme import (
    CARD_BG, FONT_BODY, FONT_HEADING,
    NAVY, TEXT_BODY, TEXT_DARK, WHITE,
)

logger = logging.getLogger(__name__)

# Alternating row colour
ROW_ALT = RGBColor(0xF0, 0xF4, 0xFF)   # very light blue


def _set_cell_fill(cell, color: RGBColor) -> None:
    """Set solid background fill on a table cell via XML."""
    try:
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        from lxml import etree
        # Remove existing fills
        for tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill')):
            for el in tcPr.findall(tag):
                tcPr.remove(el)
        solid = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb  = etree.SubElement(solid, qn('a:srgbClr'))
        srgb.set('val', f'{color.rgb:06X}')
    except Exception as exc:
        logger.debug('_set_cell_fill: %s', exc)


def _set_cell_border(cell, color: RGBColor, width_pt: float = 0.5) -> None:
    """Apply thin border on all sides of a table cell via XML."""
    try:
        from lxml import etree
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        border_xml = (
            f'<a:ln w="{int(width_pt * 12700)}" '
            f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:solidFill>'
            f'<a:srgbClr val="{color.rgb:06X}"/>'
            f'</a:solidFill></a:ln>'
        )
        for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            ln_el = etree.fromstring(border_xml.replace('<a:ln', f'<{tag}').replace('</a:ln>', f'</{tag}>'))
            existing = tcPr.find(qn(tag))
            if existing is not None:
                tcPr.remove(existing)
            tcPr.append(ln_el)
    except Exception as exc:
        logger.debug('_set_cell_border: %s', exc)


def add_table(
    slide: Slide,
    x, y, w, h,
    headers: List[str],
    rows: List[List[str]],
    caption: str = '',
) -> None:
    """
    Add a styled table to the slide.

    Args:
        slide   : target slide
        x,y,w,h : position/size in pptx units (Emu)
        headers : column header strings
        rows    : data rows (list of lists)
        caption : optional caption below the table
    """
    if not headers:
        logger.warning('add_table: no headers, skipping')
        return

    # Cap display size
    MAX_ROWS = 15
    MAX_COLS = 8
    headers = headers[:MAX_COLS]
    rows    = [row[:MAX_COLS] for row in rows[:MAX_ROWS]]

    n_rows = len(rows) + 1   # +1 for header
    n_cols = len(headers)

    try:
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
        tbl = tbl_shape.table

        # Auto font size: smaller for larger tables
        base_font  = 13 if n_rows <= 6  else (11 if n_rows <= 10 else 9)
        header_font = base_font + 1

        border_color = RGBColor(0xD1, 0xD5, 0xDB)

        # ── Header row ───────────────────────────────────────────────────────
        for col_idx, header_text in enumerate(headers):
            cell = tbl.cell(0, col_idx)
            cell.text = header_text
            _set_cell_fill(cell, NAVY)
            _set_cell_border(cell, border_color)

            tf = cell.text_frame
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.bold       = True
                    run.font.size       = Pt(header_font)
                    run.font.name       = FONT_HEADING
                    run.font.color.rgb  = WHITE

        # ── Data rows ────────────────────────────────────────────────────────
        for row_idx, row_data in enumerate(rows):
            fill_color = ROW_ALT if row_idx % 2 == 0 else WHITE
            for col_idx in range(n_cols):
                cell = tbl.cell(row_idx + 1, col_idx)
                cell_text = row_data[col_idx] if col_idx < len(row_data) else ''
                cell.text = str(cell_text)
                _set_cell_fill(cell, fill_color)
                _set_cell_border(cell, border_color)

                tf = cell.text_frame
                for para in tf.paragraphs:
                    # Right-align numbers, left-align text
                    try:
                        float(cell_text.replace(',', '').replace('%', ''))
                        para.alignment = PP_ALIGN.RIGHT
                    except (ValueError, AttributeError):
                        para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.size      = Pt(base_font)
                        run.font.name      = FONT_BODY
                        run.font.color.rgb = TEXT_BODY

    except Exception as exc:
        logger.error('add_table failed: %s', exc)
        return

    # ── Caption ──────────────────────────────────────────────────────────────
    if caption:
        from pptx.util import Inches
        cap_y = y + h + Inches(0.05)
        txb = slide.shapes.add_textbox(x, cap_y, w, Inches(0.30))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        p.text      = caption
        p.alignment = PP_ALIGN.RIGHT
        r = p.runs[0]
        r.font.size       = Pt(9)
        r.font.italic     = True
        r.font.name       = FONT_BODY
        r.font.color.rgb  = RGBColor(0x6B, 0x72, 0x80)
