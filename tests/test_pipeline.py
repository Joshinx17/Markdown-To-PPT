import unittest
from pathlib import Path
from uuid import uuid4

from click.testing import CliRunner
from pptx import Presentation

from cli import main
from orchestrator import convert
from parser.md_parser import parse_markdown
from structurer.llm_structurer import structure_presentation
from structurer.slide_types import SlideType


SAMPLE_PATH = Path(__file__).resolve().parents[1] / 'samples' / 'enterprise_ai.md'
OUTPUT_DIR = Path(__file__).resolve().parents[1] / 'output' / 'test_artifacts'
INPUT_DIR = Path(__file__).resolve().parents[1] / 'input'


class ParserTests(unittest.TestCase):
    def test_parse_markdown_extracts_sections_and_tables(self) -> None:
        doc = parse_markdown(SAMPLE_PATH.read_text(encoding='utf-8'))

        self.assertEqual(doc.title, 'The Future of Artificial Intelligence in Enterprise')
        self.assertGreaterEqual(len(doc.sections), 5)
        self.assertGreaterEqual(len(doc.tables), 2)
        self.assertTrue(doc.has_numerical_data)
        self.assertTrue(doc.has_tabular_data)


class PipelineTests(unittest.TestCase):
    def setUp(self) -> None:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        INPUT_DIR.mkdir(parents=True, exist_ok=True)
        for stale_file in OUTPUT_DIR.glob('*.pptx'):
            try:
                stale_file.unlink()
            except PermissionError:
                pass

    def unique_output(self, stem: str) -> Path:
        return OUTPUT_DIR / f'{stem}_{uuid4().hex}.pptx'

    def test_convert_succeeds_without_api_key(self) -> None:
        output_path = self.unique_output('offline_output')
        result = convert(
            input_path=str(SAMPLE_PATH),
            output_path=str(output_path),
            api_key=None,
            min_slides=10,
            max_slides=12,
            verbose=False,
        )

        self.assertTrue(result.exists())
        prs = Presentation(str(result))
        self.assertGreaterEqual(len(prs.slides), 10)

    def test_convert_uses_supplied_master_template(self) -> None:
        template_path = self.unique_output('custom_master_template')
        output_path = self.unique_output('templated_output')

        template = Presentation()
        template.slide_width = 9144000
        template.slide_height = 5143500
        template.save(str(template_path))

        result = convert(
            input_path=str(SAMPLE_PATH),
            output_path=str(output_path),
            api_key=None,
            template_path=str(template_path),
            min_slides=10,
            max_slides=12,
            verbose=False,
        )

        prs = Presentation(str(result))
        self.assertEqual(prs.slide_width, template.slide_width)
        self.assertEqual(prs.slide_height, template.slide_height)

    def test_offline_structuring_prefers_visual_slide_types(self) -> None:
        doc = parse_markdown(SAMPLE_PATH.read_text(encoding='utf-8'))
        blueprint = structure_presentation(doc, api_key=None, min_slides=10, max_slides=12)

        slide_types = {slide.type for slide in blueprint.slides}
        visual_types = {
            SlideType.EXECUTIVE_SUMMARY,
            SlideType.CHART_BAR,
            SlideType.CHART_PIE,
            SlideType.CHART_LINE,
            SlideType.PROCESS_FLOW,
            SlideType.TIMELINE,
            SlideType.COMPARISON,
            SlideType.TWO_COLUMN,
        }

        self.assertGreaterEqual(len(slide_types & visual_types), 4)
        self.assertIn(SlideType.CHART_LINE, slide_types)
        self.assertIn(SlideType.PROCESS_FLOW, slide_types)
        self.assertIn(SlideType.COMPARISON, slide_types)

    def test_cli_allows_offline_mode(self) -> None:
        runner = CliRunner()
        output_path = self.unique_output('cli_output')
        result = runner.invoke(
            main,
            [
                '--input',
                str(SAMPLE_PATH),
                '--output',
                str(output_path),
                '--slides',
                '10',
            ],
            env={'GEMINI_API_KEY': ''},
        )

        self.assertEqual(result.exit_code, 0, msg=result.output)
        self.assertIn('Offline fallback', result.output)
        self.assertTrue(output_path.exists())

    def test_cli_accepts_template_option(self) -> None:
        runner = CliRunner()
        template_path = self.unique_output('cli_master_template')
        output_path = self.unique_output('cli_template_output')

        template = Presentation()
        template.slide_width = 9144000
        template.slide_height = 5143500
        template.save(str(template_path))

        result = runner.invoke(
            main,
            [
                '--input',
                str(SAMPLE_PATH),
                '--output',
                str(output_path),
                '--slides',
                '10',
                '--template',
                str(template_path),
            ],
            env={'GEMINI_API_KEY': ''},
        )

        self.assertEqual(result.exit_code, 0, msg=result.output)
        self.assertIn(str(template_path), result.output)
        prs = Presentation(str(output_path))
        self.assertEqual(prs.slide_width, template.slide_width)

    def test_cli_can_autodiscover_input_directory_file(self) -> None:
        runner = CliRunner()
        input_copy = INPUT_DIR / 'autodiscover_test.md'
        input_copy.write_text(SAMPLE_PATH.read_text(encoding='utf-8'), encoding='utf-8')
        output_path = self.unique_output('autodiscover_output')

        try:
            result = runner.invoke(
                main,
                [
                    '--output',
                    str(output_path),
                    '--slides',
                    '10',
                ],
                env={'GEMINI_API_KEY': ''},
            )
        finally:
            try:
                input_copy.unlink(missing_ok=True)
            except PermissionError:
                pass

        self.assertEqual(result.exit_code, 0, msg=result.output)
        self.assertTrue(output_path.exists())

    def test_convert_auto_discovers_templates_directory_file(self) -> None:
        template_path = Path(__file__).resolve().parents[1] / 'templates' / '000_auto_detect_template.pptx'
        template = Presentation()
        template.slide_width = 9144000
        template.slide_height = 5143500
        template.save(str(template_path))

        output_path = self.unique_output('auto_template_output')
        try:
            result = convert(
                input_path=str(SAMPLE_PATH),
                output_path=str(output_path),
                api_key=None,
                min_slides=10,
                max_slides=12,
                verbose=False,
            )
            self.assertTrue(output_path.exists())
            prs = Presentation(str(result))
            self.assertEqual(prs.slide_width, template.slide_width)
            self.assertEqual(prs.slide_height, template.slide_height)
        finally:
            try:
                template_path.unlink(missing_ok=True)
            except PermissionError:
                pass


if __name__ == '__main__':
    unittest.main()
